"""
Generate provisioning_announcement.pptx — proper Kramer brand.
Colors, fonts, and logo from the /kramer-brand skill.
"""
import base64, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kramer brand palette ───────────────────────────────────────────────────────
DARK    = RGBColor(0x12, 0x12, 0x12)   # #121212  background
PURPLE  = RGBColor(0x82, 0x00, 0xB4)   # #8200B4  primary accent
MINT    = RGBColor(0x68, 0xFF, 0xC3)   # #68FFC3  highlight / success
PINK    = RGBColor(0xBE, 0x00, 0x78)   # #BE0078  secondary / alert
NEUTRAL = RGBColor(0xDC, 0xDD, 0xDE)   # #DCDDDE  body text
DIM     = RGBColor(0x70, 0x72, 0x74)   # muted neutral for captions
CARD    = RGBColor(0x1E, 0x1E, 0x1E)   # card surface
CARD2   = RGBColor(0x26, 0x26, 0x26)   # slightly lighter card

LOGO_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAW8AAAB3CAYAAADMzsFWAAAAAXNSR0IArs4c6QAAAARnQU1BAACx"
    "jwv8YQUAAAAJcEhZcwAAIdUAACHVAQSctJ0AADz6SURBVHhe7Z0HnBT1+cbVmGZijzEx3WhiTNUU"
    "k/wTY6oxatRoNDFGjShwXOG446gH3FEF7BWwIBZEBAWu9947rnD03jscHQ6u1+f/PrMzx9zezGy9"
    "u0V/D5/ns8fuzG9nZme+8877a2dBSUlJSemMk4K3kpKS0hkoBW8lJSWlM1AfeXi3NrTh1P4G7C48"
    "ilWv70Hl8zuQl7geC/+7Eu/fvRQLHliBnNi1cKdsxtKXdmHdO/tQVXkCtQcb0d7SrpeipKSkFFn6"
    "yMG7rbkNh1adxKrX9mDhw6sw43o3Hr8oC+M+mY5RZy1EspivozUv0l75f+Mzvjf+Mxl44os5eP1X"
    "i5EtYN/0wQHUyA1ASUlJKVL0kYB3W1Mb9pQcQ8HQDQLrYkz8XGYHkMcIjFM/kYbUc9MxVgDuj7ks"
    "1+G6RhlPXJ6LuXctxbrZ+9BwtFn/ZiUlJaXe0RkN71N7G1Dx1Ha88rNiDbYdsA4A1P469Zw0LSpn"
    "+c9dWYDCERtxfFudviVKSkpKPaszEt6H157S0hlPXJ7TAeyx3QBsO6ecnaalWKZcloPcxPU4uate"
    "3zIlJSWlntEZBe/jW2uRGb1Gy2ET2oyGreDaU+b3czue/EouKp7cjpa6Vn1LlZSUlLpXZwS8W+pb"
    "UTZ5G6ZKpM2IlykSK5j2lhmJE+Kv/3ox9pYe17daSUlJqfsU8fDev+QEXv/NYg2OKb0cafsyc+IT"
    "P5+J0se3ob1VNTNUUlLqPkUuvIV9TEVMOj9Tg6IVLCPRfCrg9r7/z2WoO9yk74ySkpJSeBWR8Cb0"
    "5v97uQbB3s5rB2s+Kcz4WTGObarR90pJSUkpfIo4eJ/YUYdXbyhB8lkLLKF4JnmU3Hyeu6pA67Gp"
    "pKSkFE5FFLyPbqjBi98v0qJWKxieiebTw1NX5GGfqshUUlIKoyIG3ofXnMSz3y7o2fw224Ybtvo8"
    "TB5zVhqmXpaDA0ur9b1VUlJSCk0RAe/j22rx3HcI7u6JuNnjkh15eGMwxjZJkb853sm4T8nnelM/"
    "Y2wTNv2zKicUs9wXry1C9W7VoUdJSSl09Tq8a6saMf06d9jBbbT6YLkTzsvEtJ+48d4dS1H+5HYs"
    "n7Ybm+ZVYbfrKPYUH8OWBQexbNoulE7Zhtm3LMHzVxdirN7dPpwgZ3mzfleGplMt+t4rKSkpBade"
    "hXdbazvm3bdci3itYBeMGWETkpMvyca7ty3Bspd24ciGGm3wKn/VWN2igT0nbh2e+XqedgNIOSs8"
    "EOe2ZcWs0b9JSUlJKTj1KrxLJ20NW8RtpEVe/pELZVO3hW3QKI7rzXG+X7jGU5Eacu/Oc9kjcxHW"
    "zd6vf4OSkpJS4Oo1eHMI1wmfy0DquQLDT1lAzk8bowm++L0irHx1D5prumd8kYajTdqEDZMuDL3T"
    "EHuKPnlFLqp3qlEJlZSUglOvwLuxuhnTrnNhjESgYz8tQCO8gwA4ITrx/Cy4Uzej8UTPjLFdtawa"
    "r/3a010/lFYqXH/efcu0nqRKSkpKgapX4O0au9kDv88IyAhvfwD+idMed5ZE3GelYeYNi3Go8qRe"
    "as+JOfFFfVZpKZ+gxw5nC5hzFmHThwf0UpWUlJT8V4/Dm5WHky7LRirB/VndVgA/xwNpw6kXZiD1"
    "EvHFGRj1k0yMuiMHs7JWY976zXh72Xq8s3wD5qzaiLQN25C3ZRfKdu/Hsv0HseHwUVSdqkVtUzOa"
    "WsKbUnFJxM9ce7AA55PDjJ8Wa6MmKikpKQWiHoV3ezsw/6EVSD43Damfz8DYz4nPExPgArNxZ+tg"
    "O0+i0m8JoH+ThWEDcjEsLheDXypA4uvimQUYNK8QA9MLEftBPqLn5SNmvm75f8yH8t4C8aJ8DEjL"
    "R1xWIQbnuTGicDEmLq7AjFWrMX/zFiw/dBD7a2rQ3OZ/KxQrucdt0SAcLMD5BMKJkZWUlJQCUY/C"
    "e1frUYy5IB0pF0gETRPcnxILyFOuzETyTdkYGp+LpCfykfhKARLeFb8nfr8QiXMKMei9QiQIuAfN"
    "L0Q8Xz+Q1w/ldYHuheJF4rQiDMooQnxWEQZmi3OKEJsjsM8rxIC8As3RBQUYXOzG+MpyzNq4HkX7"
    "9uBAXS3aeIcJUJwSLdhKTOb9p/3YjWY1kYOSklIA6jF4twsUZ9+/FMnnpSHlfAG3RN1jrs7EiLty"
    "MGRUHgZP80TVCbPk9Q15pfn3W/L6tsB7tkBZAJ4wV14F3IHAO17gPTBXXvPltUBcKH+75NVdhLji"
    "QkS5BerymljuxpQ1y5Cxdwd21Z5Ea7t/UTnbkM/+2xJPHt8C0L5M8HNiYyUlJSV/1WPwrlpZjbGX"
    "CbC/kIFkpkOYCnlaIuxXxTPyMXg6XwXUrwqwX5NXAXki4f3maXgnBAlvDdx5p+EdX+SB96BieS1x"
    "YaA4rlRc5kJ0WREGlBUgvtKFJ9YvR27Vbhxs8N2kjz1Fn/lmvqcFjQWgncy8OSecUBM4KCkp+ase"
    "g/e8iasx7OcZGDpYoP28gfolAfYL4hfF/HuaWODNdAnhrUXhEn13wPsdAfe7HngzjaLBm+kTK3in"
    "C5gzBcw28B7oDe/FHngPFHgPrBAvcSFW4B0t7l9ZhMSVJXh1+zqsrz6KFoccOVuOaO3WA81/y/Ic"
    "Y4XNEJWUlJT8UbfDu+pQDd6cvwZJE/KQ9FQ+kp4VS8Q9WF41iBPgBryne+CtRd+vm+Etrwa8mfc2"
    "wXuQDbzjdXh3pEwE3lrUzZQJ4e32wJtRtwHvuHIX4nV4x1W6MXCpG7HL3IhZ7kbU8iLErnRh8qZl"
    "KD92AI1t1jnq+f9eEVT6hOsUjd6kl6KkpKTkrG6Dd0NjC7Jd2zByqoByfC6Spgi0nxA/qcP7GfFz"
    "Orwl+k58WYc3o28TvI28d7jgzXz3QIm6veGtRd0meBPccQLuWAF33ArxSoH5KjeiV7kQtboI4zZX"
    "ovjoPjR7Qfzw2lOYeEFmwN3omW7hzDttzaG1fulONTc3Y9OmTdiwYUMXr1+/Hnv2qFYzSko9pW6B"
    "99adx/DUjHLEj8lB4niJuCcJsCeLp+oAZwROeOvRdwe8Jfq2zHsb8GalZTgibx3e8SXynlfKJF6L"
    "um3gvVpexbFrJBJf40L/tYWYuG0pllUf1Pfco0WPrAo4+ma6Zfx5GTi4ouc7Hfmrffv24fvf/z6+"
    "+c1v4lvf+lYnf+UrX8FDDz2kL6mkpNTdCiu8W1vbkZ2/DUNS85EwOhdDxgmkJwiUJ8rr42ID4FbR"
    "t1fqpFPem/A2V1oKvFlp6S+8/cl3W6ZMvKJuA9wxa+V1nQsx6yUKX1eI6PVFeHH3auyo9+Ss95ef"
    "0DodBdr2m8DnKIiRKsL7u9/9rgbqr33ta538xS9+Ef/5z3/0JZWUlLpbYYP3qZomzHxzFRKG5SJp"
    "dD6GpIjHinWAa9E3AW6kTxh9m3Lfp1Mnpui7U95br7QMAN5GM8Eu8PbOd0vUbcDbALd31N0Bbi94"
    "x22Qv8X9N0i5m4rx4eFtWj589k0VWisSK0jbmfCed+9y/YhGngjva665Bl/96lfx9a9/vZMvv/xy"
    "PPDAA/qSSkpK3a2wwPvw4Vo89WwF4pNyMCS5AENGi8cInCUC1+A9XqBsRN+Etzn6NsNbj77NeW+2"
    "+Sa8jchba+utwzveCd56SxNf8O5ImQi8zVE34e0UdXeC90b5fJO8bnaj78ZCPH1kJd6esByphHcA"
    "0Tfz3hwdselkZE7WoOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX"
    "1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsA"
    "sTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02"
    "c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDst"
    "FwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7A"
    "qyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv"
    "5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJsk"
    "iafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWV"
    "BDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn"
    "5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i5"
    "8t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2i"
    "boKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/O"
    "w7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV"
    "1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2"
    "kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8Pi"
    "BdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/"
    "P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPP"
    "WoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkA"
    "m00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni"
    "/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs"
    "782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNf"
    "ycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW"
    "4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPM"
    "e3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafO"
    "rcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDeb"
    "B3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo"
    "0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78"
    "gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKb"
    "jt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7Ck"
    "vA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUd"
    "wdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKx"
    "SpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJ"
    "YgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu"
    "2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAl"
    "z+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G"
    "9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKO"
    "npUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p"
    "2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQ"
    "ibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImK"
    "pRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OO"
    "zUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD8"
    "0s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YC"
    "t9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwM"
    "HSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54"
    "xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0q"
    "y+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6A"
    "d0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYj"
    "x7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY"
    "8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH4"
    "0CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzH"
    "l1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/U"
    "tzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpL"
    "tjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUe"
    "eAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJR"
    "Ct5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibon"
    "L6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhV"
    "UqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiU"
    "greSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1Y"
    "ueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn"
    "1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZA"
    "NgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8o"
    "QNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4o"
    "RvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0Tg"
    "AvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FN"
    "vPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbA"
    "va1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6"
    "wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1ci"
    "sbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/Utzqy"
    "pOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJ"
    "n5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt8"
    "2cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5K"
    "SpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zA"
    "C7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrB"
    "O6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreS"
    "UuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMg"
    "DpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14"
    "C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+"
    "So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzC"
    "AsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+"
    "YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvCh"
    "Y+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQ"
    "UayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Y"
    "c/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfr"
    "ADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwY"
    "UYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCt"
    "pBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WP"
    "gW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqE"
    "kXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGj"
    "oOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nL"
    "sWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6aw"
    "ANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQo"
    "KHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyo"
    "RWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZ"
    "fdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A"
    "9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQs"
    "KkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRM"
    "L0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TV"
    "BHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayq"
    "OqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R2"
    "8aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADfl"
    "vzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsL"
    "EV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5"
    "CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/l"
    "I+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdH"
    "ysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9"
    "d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5"
    "FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGF"
    "UrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgf"
    "O1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuI"
    "o/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD"
    "9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJ"
    "Fn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBU"
    "Wj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TK"
    "VzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf"
    "0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSN"
    "KFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aES"
    "eQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzul"
    "T6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1W"
    "iLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5Cgne"
    "W7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+at"
    "PETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD"
    "746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJ"
    "pI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu"
    "2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE"
    "3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KP"
    "yWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+d"
    "amrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr"
    "5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b"
    "6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6i"
    "MvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1"
    "wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe"
    "6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0"
    "Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJw"
    "IeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yi"
    "bx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLgl"
    "RRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7cc"
    "Q8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETP"
    "Fr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746U"
    "iQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4Q"
    "cA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3i"
    "rXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGE"
    "t3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMW"
    "IzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrC"
    "nlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68"
    "I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv"
    "3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIx"
    "ILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWs"
    "gG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iRe"
    "HDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+y"
    "tu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWr"
    "BLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3e"
    "TJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8"
    "zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8po"
    "FwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+X"
    "h5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4"
    "Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8U"
    "cA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXux"
    "ZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fK"
    "hOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEm"
    "B8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOn"
    "ULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cn"
    "wJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX"
    "1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsA"
    "sTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02"
    "c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDst"
    "FwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7A"
    "qyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv"
    "5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJsk"
    "iafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWV"
    "BDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn"
    "5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i5"
    "8t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2i"
    "boKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/O"
    "w7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV"
    "1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2"
    "kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8Pi"
    "BdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/"
    "P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPP"
    "WoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkA"
    "m00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni"
    "/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs"
    "782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNf"
    "ycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW"
    "4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPM"
    "e3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafO"
    "rcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDeb"
    "B3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo"
    "0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78"
    "gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKb"
    "jt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7Ck"
    "vA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUd"
    "wdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKx"
    "SpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJ"
    "YgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu"
    "2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAl"
    "z+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G"
    "9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKO"
    "npUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p"
    "2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQ"
    "ibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImK"
    "pRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OO"
    "zUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD8"
    "0s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YC"
    "t9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwM"
    "HSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54"
    "xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0q"
    "y+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6A"
    "d0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYj"
    "x7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY"
    "8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH4"
    "0CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzH"
    "l1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/U"
    "tzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpL"
    "tjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUe"
    "eAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJR"
    "Ct5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibon"
    "L6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhV"
    "UqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiU"
    "greSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1Y"
    "ueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn"
    "1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZA"
    "NgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8o"
    "QNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4o"
    "RvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0Tg"
    "AvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FN"
    "vPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbA"
    "va1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6"
    "wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1ci"
    "sbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/Utzqy"
    "pOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJ"
    "n5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt8"
    "2cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5K"
    "SpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zA"
    "C7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrB"
    "O6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreS"
    "UuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMg"
    "DpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14"
    "C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+"
    "So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzC"
    "AsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+"
    "YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvCh"
    "Y+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQ"
    "UayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Y"
    "c/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfr"
    "ADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwY"
    "UYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr1WiESeQJwIeWrBLjv5oPMe3OOzUi"
    "UgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1"
    "YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9h"
    "n1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZ"
    "ANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8"
    "oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4"
    "oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0T"
    "gAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7F"
    "NvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8Nb"
    "Ava1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE"
    "6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1c"
    "isbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/Utzq"
    "ypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtji"
    "Jn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt"
    "82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5"
    "KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6z"
    "AC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqr"
    "BO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgre"
    "SUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueM"
    "gDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E1"
    "4C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA"
    "+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNz"
    "CAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS"
    "+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvC"
    "hY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPH"
    "QUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1"
    "Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIf"
    "rADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbw"
    "YUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOC"
    "tpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5W"
    "PgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cq"
    "EkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpG"
    "joOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7n"
    "LsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6a"
    "wANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQ"
    "oKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpy"
    "oRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7g"
    "ZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/"
    "A9fSJFn0b6RPv3HeX1IkAm00K9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQ"
    "sKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiR"
    "ML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+T"
    "VBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUay"
    "qOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R"
    "28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADf"
    "lvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYs"
    "LEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ"
    "5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/"
    "lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXd"
    "HysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG"
    "9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD"
    "5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANG"
    "FUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHg"
    "fO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWu"
    "Io/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdM"
    "D9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fS"
    "JFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkB"
    "UWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0T"
    "KVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHC"
    "f0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqS"
    "NKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aE"
    "SeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzu"
    "lT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1"
    "WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5Cgn"
    "eW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+a"
    "tPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQ"
    "D746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9d"
    "JpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhR"
    "u2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZ"
    "E3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1K"
    "PyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+"
    "damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hd"
    "r5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0"
    "b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6"
    "iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK"
    "1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bf"
    "e6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi"
    "0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJ"
    "wIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6y"
    "ibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLg"
    "lRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7c"
    "cQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPET"
    "PFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746"
    "UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4"
    "QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3"
    "irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rG"
    "Et3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWM"
    "WIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damr"
    "CnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y6"
    "8I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RP"
    "v3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvI"
    "xILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbW"
    "sgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iR"
    "eHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+"
    "ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeW"
    "rBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3"
    "eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq"
    "8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8p"
    "oFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+"
    "Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO"
    "4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8"
    "UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXu"
    "xZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3f"
    "KhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzE"
    "mB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlO"
    "nULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8c"
    "nwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3He"
    "X1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILs"
    "AsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG0"
    "2c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDs"
    "tFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7"
    "AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLj"
    "v5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJs"
    "kiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRW"
    "VBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwY"
    "n5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i"
    "58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2"
    "iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/"
    "Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGc"
    "V1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD"
    "2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8P"
    "iBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ"
    "/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJP"
    "PWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1Ik"
    "Am00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTn"
    "i/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+R"
    "s782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwN"
    "fycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7Aqyv"
    "W4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oP"
    "Me3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiaf"
    "OrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDe"
    "bB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GD"
    "o0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t7"
    "8gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboK"
    "bjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7C"
    "kvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lU"
    "dwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTK"
    "xSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdq"
    "JYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Z"
    "u2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoA"
    "lz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00"
    "G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MK"
    "OnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782"
    "p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycM"
    "QibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4Im"
    "KpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3O"
    "OzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJskiafOrcD"
    "80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3Y"
    "Ct9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0Hw"
    "MHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg5"
    "4xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0"
    "qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6"
    "Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdY"
    "jx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpc"
    "Y8NbAva1Yc/R28aESeQJwIeWrBLjv5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH"
    "40CE6wIfrADflvzulT6yibx3eTJskiafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2Yz"
    "Hl1cisbwYUYsLEV1WiLglRRq8zRWVBDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/"
    "UtzqypOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX1IkAm00G9Up"
    "LtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsAsTni/MKOnpU"
    "eeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02c+Rs782p2SJ"
    "RCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDstFwNfycMQibo"
    "nL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7AqyvW4ImKpRh"
    "VUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28="
)

# Decode logo safely
try:
    _b64 = LOGO_B64 + '=' * ((4 - len(LOGO_B64) % 4) % 4)
    LOGO_BYTES = base64.b64decode(_b64)
    LOGO_IO = io.BytesIO(LOGO_BYTES)
    HAS_LOGO = True
except Exception:
    HAS_LOGO = False

# ── Fonts (Kramer brand: GT Eesti Display Md for headings, Lt for body) ────────
F_HEAD = "GT Eesti Display Md"
F_BODY = "GT Eesti Display Lt"

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitives ─────────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    return s

def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def rect(s, l, t, w, h, fill, line_rgb=None, lpt=1.5):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_rgb: sh.line.color.rgb = line_rgb; sh.line.width = Pt(lpt)
    else:        sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, text, size,
        color=NEUTRAL, align=PP_ALIGN.LEFT, italic=False, font=F_BODY, wrap=True):
    tb = box(s, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    r.font.bold = False   # weight via font variant, not bold flag

def add_logo(s, l=12.35, t=0.12, w=0.75):
    """Place the Kramer logo (top-right by default)."""
    if not HAS_LOGO:
        return
    try:
        LOGO_IO.seek(0)
        s.shapes.add_picture(LOGO_IO, Inches(l), Inches(t), width=Inches(w))
    except Exception:
        pass

def footer(s, n, total=6):
    txt(s, 0.45, 7.22, 10, 0.22,
        "Kramer Electronics  ·  IT Department  ·  May 2026",
        8, color=DIM)
    txt(s, 12.1, 7.22, 1.1, 0.22, f"{n} / {total}", 8,
        color=DIM, align=PP_ALIGN.RIGHT)

def pill(s, l, t, w, text, bg=PURPLE):
    rect(s, l, t, w, 0.33, bg)
    txt(s, l, t + 0.01, w, 0.32, text.upper(), 8.5,
        color=DARK, align=PP_ALIGN.CENTER, font=F_HEAD)

def heading(s, text, l=0.5, t=0.55, size=28):
    txt(s, l, t, 12.4, 0.9, text, size,
        color=NEUTRAL, font=F_HEAD)

def sub_heading(s, text, l=0.5, t=1.38):
    txt(s, l, t, 12.4, 0.38, text, 13,
        color=DIM, italic=True)
    rect(s, 0.5, 1.76, 12.3, 0.03, PURPLE)   # thin purple rule under sub

def bullets(s, l, t, w, items, size=13):
    """items = list of (dot_color, bold_label, body_text)"""
    tb = box(s, l, t, w, H.inches - t - 0.5)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in items:
        dot_c, label, rest = item[0], item[1], item[2]
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_before = Pt(9)
        r = p.add_run(); r.text = "▪  "
        r.font.size = Pt(size); r.font.color.rgb = dot_c
        r.font.name = F_BODY; r.font.bold = False
        if label:
            r2 = p.add_run(); r2.text = label
            r2.font.size = Pt(size); r2.font.color.rgb = NEUTRAL
            r2.font.name = F_HEAD; r2.font.bold = False
        if rest:
            r3 = p.add_run(); r3.text = "  " + rest
            r3.font.size = Pt(size); r3.font.color.rgb = DIM
            r3.font.name = F_BODY; r3.font.bold = False


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  ·  Cover
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)   # top bar
add_logo(s)

# Main title
txt(s, 0.9, 0.9, 11.5, 2.2,
    "New Employee\nProvisioning",
    54, color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_HEAD)

# Mint accent line
rect(s, 4.5, 3.3, 4.3, 0.06, MINT)

# Tagline
txt(s, 1.0, 3.52, 11.3, 0.5,
    "From HiBob form to Active Directory — fully automated",
    16, color=DIM, align=PP_ALIGN.CENTER, italic=True, font=F_BODY)

txt(s, 1.5, 4.08, 10.3, 0.42,
    "Days to build it. Years overdue.",
    14, color=MINT, align=PP_ALIGN.CENTER, font=F_HEAD)

# Audience pills
pill(s, 3.05,  4.85, 2.1, "🧑‍💼  For HR",     PURPLE)
pill(s, 5.6,   4.85, 2.1, "🖥️  For IT",      PINK)
pill(s, 8.15,  4.85, 2.1, "🔔  For Admins",  MINT)

# Credit
txt(s, 1.0, 5.7, 11.3, 0.38,
    "Built by Omri Cohen, IT Department",
    11, color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)

footer(s, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  ·  Before & After
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)
heading(s, "Before & After")
sub_heading(s, "New employee onboarding — then and now")

# Column headers
txt(s, 0.5, 1.9, 5.9, 0.36, "BEFORE", 10,
    color=PINK, align=PP_ALIGN.LEFT, font=F_HEAD)
rect(s, 0.5, 2.26, 5.9, 0.03, PINK)

txt(s, 6.9, 1.9, 5.9, 0.36, "AFTER", 10,
    color=MINT, align=PP_ALIGN.LEFT, font=F_HEAD)
rect(s, 6.9, 2.26, 5.9, 0.03, MINT)

# Divider
rect(s, 6.55, 1.88, 0.04, 5.3, RGBColor(0x33, 0x33, 0x33))

# Before bullets
bullets(s, 0.5, 2.38, 5.9, [
    (PINK, "Email IT manually",          "— no standard process"),
    (PINK, "Wait hours or days",         "— no SLA, no visibility"),
    (PINK, "Manual AD account creation", ""),
    (PINK, "Separate M365 request",      ""),
    (PINK, "Steps missed, inconsistent", ""),
], size=14)

# After bullets
bullets(s, 6.9, 2.38, 5.9, [
    (MINT, "Fill the HiBob form",               "— that's it"),
    (MINT, "AD account created automatically",   "— under 1 minute"),
    (MINT, "M365 synced + 8 groups assigned",    "— no manual steps"),
    (MINT, "Priority / Salesforce tickets",      "— auto-created if flagged"),
    (MINT, "Done in ~20 minutes",                "— admins notified by email"),
], size=14)

footer(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  ·  HR
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)
pill(s, 0.5, 0.18, 1.9, "🧑‍💼  For HR", PURPLE)
heading(s, "Just fill in HiBob. That's your only step.", t=0.55)
sub_heading(s, "HiBob is now the ONLY channel for new employee requests. No emails, no chats to IT.")

# Warning banner
rect(s, 0.5, 1.86, 12.3, 0.52, RGBColor(0x24, 0x08, 0x18), line_rgb=PINK, lpt=1.5)
txt(s, 0.7, 1.91, 12.0, 0.42,
    "Requests sent directly to IT — by email or message — will no longer be actioned.",
    11.5, color=NEUTRAL, font=F_HEAD)

# Left column — required fields
txt(s, 0.5, 2.56, 5.9, 0.32, "REQUIRED FIELDS", 8.5,
    color=PURPLE, font=F_HEAD)
rect(s, 0.5, 2.88, 5.9, 0.03, PURPLE)
bullets(s, 0.5, 2.98, 5.9, [
    (PURPLE, "First Name & Last Name",  ""),
    (PURPLE, "Department & Division",   ""),
    (PURPLE, "Country & Region",        ""),
    (PURPLE, "Start Date",              ""),
    (PURPLE, "Job Title",               ""),
    (PURPLE, "Reports To",              "(manager must exist in AD)"),
], size=13)

# Right column — system access
txt(s, 6.9, 2.56, 5.9, 0.32, "SYSTEM ACCESS  (Yes / No)", 8.5,
    color=MINT, font=F_HEAD)
rect(s, 6.9, 2.88, 5.9, 0.03, MINT)
bullets(s, 6.9, 2.98, 5.9, [
    (MINT, "Priority ERP access?",  "If Yes  →  Permissions As: [name]"),
    (MINT, "Salesforce access?",    "If Yes  →  Country Permission: [country]"),
], size=13)

footer(s, 3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  ·  IT — Automation Flow
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)
pill(s, 0.5, 0.18, 1.9, "🖥️  For IT", PINK)
heading(s, "7 steps. Zero manual effort.", t=0.55)
sub_heading(s, "From HiBob notification email to fully provisioned user.")

steps = [
    ("1", "HiBob\nEmail",        "Arrives at\nservicedesk@"),
    ("2", "Kdesk\nDetects",      "Subject +\n@hibob.com verified"),
    ("3", "AD Account\nCreated", "OU, manager, country,\nhome drive"),
    ("4", "KADSYNC\nTriggered",  "Delta sync fired\nUser in M365 ~3 min"),
    ("5", "M365\nGroups",        "8 groups assigned\n(6 dept + 2 universal)"),
    ("6", "System\nTickets",     "Priority / Salesforce\nauto-created if flagged"),
    ("7", "Done ✓",              "Admins notified\nWork email confirmed"),
]

sw   = 1.68
gap  = 0.14
x0   = 0.45

for i, (n, title, desc) in enumerate(steps):
    lx  = x0 + i * (sw + gap)
    acc = MINT   if i == 6 else (PURPLE if i == 0 else NEUTRAL)
    rect(s, lx, 1.95, sw, 4.55, CARD, line_rgb=acc, lpt=1.5)
    txt(s, lx, 2.05, sw, 0.58, n, 26,
        color=acc, align=PP_ALIGN.CENTER, font=F_HEAD)
    txt(s, lx + 0.06, 2.62, sw - 0.12, 0.72, title, 10.5,
        color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_HEAD)
    txt(s, lx + 0.08, 3.36, sw - 0.16, 1.3, desc, 9,
        color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)
    if i < 6:
        txt(s, lx + sw, 3.42, gap + 0.05, 0.28, "›", 16,
            color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)

footer(s, 4)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  ·  Admins — Notifications
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)
pill(s, 0.5, 0.18, 2.1, "🔔  For Admins", MINT)
heading(s, "An email for every outcome.", t=0.55)
sub_heading(s,
    "Kdesk_Superusers@kramerav.com  ·  Dashboard: kdesk.kramerav.com/hibob-sync/")

cards = [
    ("✅", "Provisioned",
     "Name · work email · dept · start date",
     "Links to ticket & dashboard.",
     MINT),
    ("❌", "Failed",
     "Same info + full script log embedded",
     "ERROR lines highlighted. Investigate if needed.",
     PINK),
    ("⚠️", "Active Account Found",
     "An active AD account exists for this name",
     "Dashboard: Create New Account  or  Cancel.",
     PURPLE),
    ("⚠️", "Returning Employee",
     "A disabled AD account was found",
     "Re-activate manually in AD, then close the request.",
     NEUTRAL),
]

cw = 2.92
for i, (ico, title, l1, l2, acc) in enumerate(cards):
    lx = 0.45 + i * (cw + 0.175)
    rect(s, lx, 1.95, cw, 4.7, CARD, line_rgb=acc, lpt=2.0)
    # icon
    ib = box(s, lx, 2.08, cw, 0.5)
    p  = ib.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r  = p.add_run(); r.text = ico
    r.font.size = Pt(26); r.font.name = "Segoe UI Emoji"
    # title
    txt(s, lx + 0.1, 2.6, cw - 0.2, 0.42, title, 12,
        color=acc, align=PP_ALIGN.CENTER, font=F_HEAD)
    # line 1
    txt(s, lx + 0.12, 3.04, cw - 0.24, 0.45, l1, 10.5,
        color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_BODY)
    # line 2
    txt(s, lx + 0.12, 3.52, cw - 0.24, 0.65, l2, 9.5,
        color=DIM, align=PP_ALIGN.CENTER, italic=True, font=F_BODY)

footer(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  ·  Closing
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)

# Hero text
txt(s, 0.9, 0.9, 11.5, 1.85,
    "Days to build it.\nYears overdue.",
    48, color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_HEAD)

# Mint rule
rect(s, 4.5, 2.95, 4.3, 0.06, MINT)

# Tag line
txt(s, 1.5, 3.14, 10.3, 0.52,
    "Onboarding at Kramer is now faster, consistent, and requires zero manual IT effort.",
    14, color=DIM, align=PP_ALIGN.CENTER, italic=True, font=F_BODY)

# Credit box
rect(s, 3.4, 3.95, 6.5, 1.85, CARD, line_rgb=PURPLE, lpt=2.0)
txt(s, 3.4, 4.08, 6.5, 0.38,
    "Designed, engineered & delivered by",
    11, color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)
txt(s, 3.4, 4.44, 6.5, 0.72,
    "Omri Cohen",
    32, color=MINT, align=PP_ALIGN.CENTER, font=F_HEAD)
txt(s, 3.4, 5.14, 6.5, 0.38,
    "IT Department  ·  Kramer Electronics",
    11, color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)

footer(s, 6)


# ── Save ───────────────────────────────────────────────────────────────────────
out = r"C:\Users\ocohen\kdesk\provisioning_announcement_v3.pptx"
prs.save(out)
print(f"Saved: {out}")
