"""
Generate termination_slide.pptx — single slide, Kramer brand.
Same design as provisioning_announcement.pptx.
"""
import base64, io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kramer brand palette ───────────────────────────────────────────────────────
DARK    = RGBColor(0x12, 0x12, 0x12)
PURPLE  = RGBColor(0x82, 0x00, 0xB4)
MINT    = RGBColor(0x68, 0xFF, 0xC3)
PINK    = RGBColor(0xBE, 0x00, 0x78)
NEUTRAL = RGBColor(0xDC, 0xDD, 0xDE)
DIM     = RGBColor(0x70, 0x72, 0x74)
CARD    = RGBColor(0x1E, 0x1E, 0x1E)

F_HEAD = "GT Eesti Display Md"
F_BODY = "GT Eesti Display Lt"

LOGO_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAW8AAAB3CAYAAADMzsFWAAAAAXNSR0IArs4c6QAAAARnQU1BAACx"
    "jwv8YQUAAAAJcEhZcwAAIdUAACHVAQSctJ0AADz6SURBVHhe7Z0HnBT1+cbVmGZijzEx3WhiTNUU"
    "k/wTY6oxatRoNDFGjShwXOG446gH3FEF7BWwIBZEBAWu9947rnD03jscHQ6u1+f/PrMzx9zezGy9"
    "u0V/D5/ns8fuzG9nZme+8877a2dBSUlJSemMk4K3kpKS0hkoBW8lJSWlM1AfeXi3NrTh1P4G7C48"
    "ilWv70Hl8zuQl7geC/+7Eu/fvRQLHliBnNi1cKdsxtKXdmHdO/tQVXkCtQcb0d7SrpeipKSkFFn6"
    "yMG7rbkNh1adxKrX9mDhw6sw43o3Hr8oC+M+mY5RZy1EspivozUv0l75f+Mzvjf+Mxl44os5eP1X"
    "i5EtYN/0wQHUyA1ASUlJKVL0kYB3W1Mb9pQcQ8HQDQLrYkz8XGYHkMcIjFM/kYbUc9MxVgDuj7ks"
    "1+G6RhlPXJ6LuXctxbrZ+9BwtFn/ZiUlJaXe0RkN71N7G1Dx1Ha88rNiDbYdsA4A1P469Zw0LSpn"
    "+c9dWYDCERtxfFudviVKSkpKPaszEt6H157S0hlPXJ7TAeyx3QBsO6ecnaalWKZcloPcxPU4uate"
    "3zIlJSWlntEZBe/jW2uRGb1Gy2ET2oyGreDaU+b3czue/EouKp7cjpa6Vn1LlZSUlLpXZwS8W+pb"
    "UTZ5G6ZKpM2IlykSK5j2lhmJE+Kv/3ox9pYe17daSUlJqfsU8fDev+QEXv/NYg2OKb0cafsyc+IT"
    "P5+J0se3ob1VNTNUUlLqPkUuvIV9TEVMOj9Tg6IVLCPRfCrg9r7/z2WoO9yk74ySkpJSeBWR8Cb0"
    "5v97uQbB3s5rB2s+Kcz4WTGObarR90pJSUkpfIo4eJ/YUYdXbyhB8lkLLKF4JnmU3Hyeu6pA67Gp"
    "pKSkFE5FFLyPbqjBi98v0qJWKxieiebTw1NX5GGfqshUUlIKoyIG3ofXnMSz3y7o2fw224Ybtvo8"
    "TB5zVhqmXpaDA0ur9b1VUlJSCk0RAe/j22rx3HcI7u6JuNnjkh15eGMwxjZJkb853sm4T8nnelM/"
    "Y2wTNv2zKicUs9wXry1C9W7VoUdJSSl09Tq8a6saMf06d9jBbbT6YLkTzsvEtJ+48d4dS1H+5HYs"
    "n7Ybm+ZVYbfrKPYUH8OWBQexbNoulE7Zhtm3LMHzVxdirN7dPpwgZ3mzfleGplMt+t4rKSkpBade"
    "hXdbazvm3bdci3itYBeMGWETkpMvyca7ty3Bspd24ciGGm3wKn/VWN2igT0nbh2e+XqedgNIOSs8"
    "EOe2ZcWs0b9JSUlJKTj1KrxLJ20NW8RtpEVe/pELZVO3hW3QKI7rzXG+X7jGU5Eacu/Oc9kjcxHW"
    "zd6vf4OSkpJS4Oo1eHMI1wmfy0DquQLDT1lAzk8bowm++L0irHx1D5prumd8kYajTdqEDZMuDL3T"
    "EHuKPnlFLqp3qlEJlZSUglOvwLuxuhnTrnNhjESgYz8tQCO8gwA4ITrx/Cy4Uzej8UTPjLFdtawa"
    "r/3a010/lFYqXH/efcu0nqRKSkpKgapX4O0au9kDv88IyAhvfwD+idMed5ZE3GelYeYNi3Go8qRe"
    "as+JOfFFfVZpKZ+gxw5nC5hzFmHThwf0UpWUlJT8V4/Dm5WHky7LRirB/VndVgA/xwNpw6kXZiD1"
    "EvHFGRj1k0yMuiMHs7JWY976zXh72Xq8s3wD5qzaiLQN25C3ZRfKdu/Hsv0HseHwUVSdqkVtUzOa"
    "WsKbUnFJxM9ce7AA55PDjJ8Wa6MmKikpKQWiHoV3ezsw/6EVSD43Damfz8DYz4nPExPgArNxZ+tg"
    "O0+i0m8JoH+ThWEDcjEsLheDXypA4uvimQUYNK8QA9MLEftBPqLn5SNmvm75f8yH8t4C8aJ8DEjL"
    "R1xWIQbnuTGicDEmLq7AjFWrMX/zFiw/dBD7a2rQ3OZ/KxQrucdt0SAcLMD5BMKJkZWUlJQCUY/C"
    "e1frUYy5IB0pF0gETRPcnxILyFOuzETyTdkYGp+LpCfykfhKARLeFb8nfr8QiXMKMei9QiQIuAfN"
    "L0Q8Xz+Q1w/ldYHuheJF4rQiDMooQnxWEQZmi3OKEJsjsM8rxIC8As3RBQUYXOzG+MpyzNq4HkX7"
    "9uBAXS3aeIcJUJwSLdhKTOb9p/3YjWY1kYOSklIA6jF4twsUZ9+/FMnnpSHlfAG3RN1jrs7EiLty"
    "MGRUHgZP80TVCbPk9Q15pfn3W/L6tsB7tkBZAJ4wV14F3IHAO17gPTBXXvPltUBcKH+75NVdhLji"
    "QkS5BerymljuxpQ1y5Cxdwd21Z5Ea7t/UTnbkM/+2xJPHt8C0L5M8HNiYyUlJSV/1WPwrlpZjbGX"
    "CbC/kIFkpkOYCnlaIuxXxTPyMXg6XwXUrwqwX5NXAXki4f3maXgnBAlvDdx5p+EdX+SB96BieS1x"
    "YaA4rlRc5kJ0WREGlBUgvtKFJ9YvR27Vbhxs8N2kjz1Fn/lmvqcFjQWgncy8OSecUBM4KCkp+ase"
    "g/e8iasx7OcZGDpYoP28gfolAfYL4hfF/HuaWODNdAnhrUXhEn13wPsdAfe7HngzjaLBm+kTK3in"
    "C5gzBcw28B7oDe/FHngPFHgPrBAvcSFW4B0t7l9ZhMSVJXh1+zqsrz6KFoccOVuOaO3WA81/y/Ic"
    "Y4XNEJWUlJT8UbfDu+pQDd6cvwZJE/KQ9FQ+kp4VS8Q9WF41iBPgBryne+CtRd+vm+Etrwa8mfc2"
    "wXuQDbzjdXh3pEwE3lrUzZQJ4e32wJtRtwHvuHIX4nV4x1W6MXCpG7HL3IhZ7kbU8iLErnRh8qZl"
    "KD92AI1t1jnq+f9eEVT6hOsUjd6kl6KkpKTkrG6Dd0NjC7Jd2zByqoByfC6Spgi0nxA/qcP7GfFz"
    "Orwl+k58WYc3o28TvI28d7jgzXz3QIm6veGtRd0meBPccQLuWAF33ArxSoH5KjeiV7kQtboI4zZX"
    "ovjoPjR7Qfzw2lOYeEFmwN3omW7hzDttzaG1fulONTc3Y9OmTdiwYUMXr1+/Hnv2qFYzSko9pW6B"
    "99adx/DUjHLEj8lB4niJuCcJsCeLp+oAZwROeOvRdwe8Jfq2zHsb8GalZTgibx3e8SXynlfKJF6L"
    "um3gvVpexbFrJBJf40L/tYWYuG0pllUf1Pfco0WPrAo4+ma6Zfx5GTi4ouc7Hfmrffv24fvf/z6+"
    "+c1v4lvf+lYnf+UrX8FDDz2kL6mkpNTdCiu8W1vbkZ2/DUNS85EwOhdDxgmkJwiUJ8rr42ID4FbR"
    "t1fqpFPem/A2V1oKvFlp6S+8/cl3W6ZMvKJuA9wxa+V1nQsx6yUKX1eI6PVFeHH3auyo9+Ss95ef"
    "0DodBdr2m8DnKIiRKsL7u9/9rgbqr33ta508xS9+Ef/5z3/0JZWUlLpbYYP3qZomzHxzFRKG5SJp"
    "dD6GpIjHinWAa9E3AW6kTxh9m3Lfp1Mnpui7U95br7QMAN5GM8Eu8PbOd0vUbcDbALd31N0Bbi94"
    "x22Qv8X9N0i5m4rx4eFtWj589k0VWisSK0jbmfCed+9y/YhGngjva665Bl/96lfx9a9/vZMvv/xy"
    "PPDAA/qSSkpK3a2wwPvw4Vo89WwF4pNyMCS5AENGi8cInCUC1+A9XqBsRN+Etzn6NsNbj77NeW+2"
    "+Sa8jchba+utwzveCd56SxNf8O5ImQi8zVE34e0UdXeC90b5fJO8bnaj78ZCPH1kJd6esByphHcA"
    "0Tfz3hwdselkZE7WoOCtpBQ5CgneW7ccQ8poFwYn5GDo0HwMHSZANgA+So/A9fSJFn0b6RPv3HeX"
    "1IkAm00G9UpLtjiJn5WPgW/lI+atPETPFr+Xh5i58t78gg54xy8oQNzCAsQsKkBUWj6iMvIxILsA"
    "sTni/MKOnpUeeAt82cqEkXdHysQD746UiQO4Yy2iboKbjt0qy+4oRvS+YiRML0TKVzK1wbWsgG02"
    "c+Rs782p2SJRCt5KSpGjoOG9d9dJpI4QcA8UcA/Ow7CkvA6Ad0TgAvChY+TVBHCf0bfe6iReHDst"
    "FwNfycMQibonL6zAC7nLsWD5FhRu2o3irXuxZGcV1lUdwdYjx7FNvPHQUayqOqSNKFi0Yw+ytu7A"
    "qyvW4ImKpRhVUqrBO6awANGFUrZE3rGEt3fKhOD2kTKxSpcY8NbAva1Yc/R28aESeQJwIeWrBLjv"
    "5oPMe3OOzUiUgreSUuQoKHgfO1KPyWMWIzEmB8PiBdqJYgH40CE6wIfrADflvzulT6yibx3eTJsk"
    "iafOrcD80s1YueMgDpyoRWuIo/+damrCnlOnULZ/P+Zu2YzHl1cisbwYUYsLEV1WiLglRRq8zRWV"
    "BDebB3YCt9hn1E14C7gZfdMD9hdr5Y68I8cnwJPPWoAlz+/UtzqypOA="
)

try:
    _b64 = LOGO_B64 + '=' * ((4 - len(LOGO_B64) % 4) % 4)
    LOGO_BYTES = base64.b64decode(_b64)
    LOGO_IO    = io.BytesIO(LOGO_BYTES)
    HAS_LOGO   = True
except Exception:
    HAS_LOGO = False

# ── Canvas ─────────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Primitives ─────────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    return s

def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def rect(s, l, t, w, h, fill, line_rgb=None, lpt=1.5):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_rgb: sh.line.color.rgb = line_rgb; sh.line.width = Pt(lpt)
    else:        sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, text, size,
        color=NEUTRAL, align=PP_ALIGN.LEFT, italic=False, font=F_BODY, wrap=True):
    tb = box(s, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    r.font.bold = False

def add_logo(s):
    if not HAS_LOGO:
        return
    try:
        LOGO_IO.seek(0)
        s.shapes.add_picture(LOGO_IO, Inches(12.35), Inches(0.12), width=Inches(0.75))
    except Exception:
        pass

def pill(s, l, t, w, text, bg=PURPLE):
    rect(s, l, t, w, 0.33, bg)
    txt(s, l, t + 0.01, w, 0.32, text.upper(), 8.5,
        color=DARK, align=PP_ALIGN.CENTER, font=F_HEAD)

def heading(s, text, l=0.5, t=0.55, size=28):
    txt(s, l, t, 12.4, 0.9, text, size, color=NEUTRAL, font=F_HEAD)

def sub_heading(s, text, l=0.5, t=1.38):
    txt(s, l, t, 12.4, 0.38, text, 13, color=DIM, italic=True)
    rect(s, 0.5, 1.76, 12.3, 0.03, PURPLE)

def footer(s):
    txt(s, 0.45, 7.22, 10, 0.22,
        "Kramer Electronics  ·  IT Department  ·  May 2026",
        8, color=DIM)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE  ·  Offboarding — Also Automated
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)
pill(s, 0.5, 0.18, 2.4, "🎉  Now Also Automated", MINT)
heading(s, "Offboarding is now fully automated too.", t=0.55)
sub_heading(s, "From HiBob termination email to fully offboarded — zero manual IT steps.")

steps = [
    ("1", "HiBob\nEmail",        "Termination email\narrives at\nservicedesk@"),
    ("2", "Scheduled\n23:59",    "Waits until end of\nemployee's last day\nlocal time"),
    ("3", "AD\nDisabled",        "Account disabled,\ngroups stripped,\nmoved to deletion OU"),
    ("4", "Mailbox\n→ Shared",   "Converted to Shared.\nManager gets\nFull Access"),
    ("5", "M365\nGroups",        "Removed from all\nAAD + EXO\ndistribution groups"),
    ("6", "OneDrive\nAccess",    "Manager granted\nSite Collection Admin\non employee's drive"),
    ("7", "Done ✓",              "Priority + Salesforce\ntermination tickets\ncreated automatically"),
]

sw  = 1.68
gap = 0.14
x0  = 0.45

for i, (n, title, desc) in enumerate(steps):
    lx  = x0 + i * (sw + gap)
    acc = MINT if i == 6 else (PURPLE if i == 0 else NEUTRAL)
    rect(s, lx, 1.95, sw, 4.55, CARD, line_rgb=acc, lpt=1.5)
    txt(s, lx, 2.05, sw, 0.58, n, 26,
        color=acc, align=PP_ALIGN.CENTER, font=F_HEAD)
    txt(s, lx + 0.06, 2.62, sw - 0.12, 0.72, title, 10.5,
        color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_HEAD)
    txt(s, lx + 0.08, 3.36, sw - 0.16, 1.3, desc, 9,
        color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)
    if i < 6:
        txt(s, lx + sw, 3.42, gap + 0.05, 0.28, "›", 16,
            color=DIM, align=PP_ALIGN.CENTER, font=F_BODY)

footer(s)

# ── Save ───────────────────────────────────────────────────────────────────────
out = r"C:\Users\ocohen\kdesk\termination_slide.pptx"
prs.save(out)
print(f"Saved: {out}")
