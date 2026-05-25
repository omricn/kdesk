# Generate a single-slide PPTX for SF & Priority Admins.
# Output: provisioning_sf_priority_slide.pptx
import base64, io, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Kramer brand palette ───────────────────────────────────────────────────────
DARK    = RGBColor(0x12, 0x12, 0x12)
PURPLE  = RGBColor(0x82, 0x00, 0xB4)
MINT    = RGBColor(0x68, 0xFF, 0xC3)
PINK    = RGBColor(0xBE, 0x00, 0x78)
NEUTRAL = RGBColor(0xDC, 0xDD, 0xDE)
DIM     = RGBColor(0x70, 0x72, 0x74)
CARD    = RGBColor(0x1E, 0x1E, 0x1E)

F_HEAD = "GT Eesti Display Md"
F_BODY = "GT Eesti Display Lt"

W = Inches(13.33)
H = Inches(7.5)

# ── Logo ───────────────────────────────────────────────────────────────────────
try:
    with open(r'C:\Users\ocohen\.claude\commands\kramer-brand.md', 'r', encoding='utf-8') as _f:
        _skill = _f.read()
    _m = re.search(r'iVBORw0KGgo[A-Za-z0-9+/=\n\r]+', _skill)
    _b64 = _m.group(0).replace('\n','').replace('\r','').strip()
    _b64 += '=' * ((4 - len(_b64) % 4) % 4)
    LOGO_IO  = io.BytesIO(base64.b64decode(_b64))
    HAS_LOGO = True
except Exception:
    HAS_LOGO = False

# ── Primitives ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

def new_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = DARK
    return s

def box(s, l, t, w, h):
    return s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def rect(s, l, t, w, h, fill, line_rgb=None, lpt=1.5):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_rgb: sh.line.color.rgb = line_rgb; sh.line.width = Pt(lpt)
    else:        sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, text, size,
        color=NEUTRAL, align=PP_ALIGN.LEFT, italic=False, font=F_BODY, wrap=True):
    tb = box(s, l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    r.font.bold = False

def pill(s, l, t, w, text, bg=PURPLE):
    rect(s, l, t, w, 0.33, bg)
    txt(s, l, t + 0.01, w, 0.32, text.upper(), 8.5,
        color=DARK, align=PP_ALIGN.CENTER, font=F_HEAD)

def add_logo(s):
    if not HAS_LOGO:
        return
    try:
        LOGO_IO.seek(0)
        s.shapes.add_picture(LOGO_IO, Inches(12.35), Inches(0.12), width=Inches(0.75))
    except Exception:
        pass

def footer(s, page_label=""):
    txt(s, 0.45, 7.22, 10, 0.22,
        "Kramer Electronics  ·  IT Department  ·  May 2026",
        8, color=DIM)
    if page_label:
        txt(s, 12.1, 7.22, 1.1, 0.22, page_label, 8,
            color=DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE  ·  SF & Priority Admins
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, 0, 0, 13.33, 0.08, PURPLE)
add_logo(s)

pill(s, 0.5, 0.18, 2.75, "🎫  For SF & Priority Admins", PURPLE)

# Heading
txt(s, 0.5, 0.62, 12.3, 0.75,
    "You'll get a ticket. It comes from the new employee.",
    30, color=NEUTRAL, font=F_HEAD)

# Sub rule + subtitle
rect(s, 0.5, 1.36, 12.3, 0.03, PURPLE)
txt(s, 0.5, 1.44, 12.3, 0.35,
    "When a new hire requires Salesforce or Priority access, Kdesk creates the ticket automatically "
    "— no manual request needed from HR or IT.",
    12, color=DIM, italic=True, font=F_BODY)

# ── Two cards ─────────────────────────────────────────────────────────────────
CW = 5.8   # card width
CX = [0.5, 6.9]  # card left positions (with center gap)

card_data = [
    ("☁️",  "Salesforce",
     "If the HiBob form had  \"Salesforce: Yes\"",
     "Country Permission specified in the form.",
     MINT),
    ("📦", "Priority ERP",
     "If the HiBob form had  \"Priority: Yes\"",
     "Permissions copied from the named template user.",
     PURPLE),
]

for (ico, title, trigger, detail, acc), cx in zip(card_data, CX):
    rect(s, cx, 1.92, CW, 3.9, CARD, line_rgb=acc, lpt=2.0)

    # Icon
    ib = box(s, cx, 2.05, CW, 0.55)
    p  = ib.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r  = p.add_run(); r.text = ico
    r.font.size = Pt(30); r.font.name = "Segoe UI Emoji"

    # Title
    txt(s, cx + 0.1, 2.62, CW - 0.2, 0.42, title, 15,
        color=acc, align=PP_ALIGN.CENTER, font=F_HEAD)

    # Trigger line
    txt(s, cx + 0.15, 3.1, CW - 0.3, 0.38, trigger, 11,
        color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_BODY)

    # Detail line
    txt(s, cx + 0.15, 3.5, CW - 0.3, 0.38, detail, 10.5,
        color=DIM, align=PP_ALIGN.CENTER, italic=True, font=F_BODY)

    # Divider
    rect(s, cx + 0.2, 3.96, CW - 0.4, 0.02, acc)

    # "Sent by" label + value
    txt(s, cx + 0.15, 4.04, CW - 0.3, 0.3, "TICKET SENDER", 8,
        color=DIM, align=PP_ALIGN.CENTER, font=F_HEAD)
    txt(s, cx + 0.15, 4.34, CW - 0.3, 0.55,
        "The new employee's work email\n(auto-created during provisioning)",
        10.5, color=NEUTRAL, align=PP_ALIGN.CENTER, font=F_BODY)

# ── Bottom note ───────────────────────────────────────────────────────────────
rect(s, 0.5, 6.04, 12.3, 0.52, RGBColor(0x1A, 0x0A, 0x22), line_rgb=PURPLE, lpt=1.2)
txt(s, 0.65, 6.1, 12.0, 0.42,
    "Tickets arrive from the new hire's address so replies go directly to them. "
    "No need to chase HR or IT — just action the ticket as usual.",
    11, color=NEUTRAL, font=F_BODY)

footer(s)

# ── Save ───────────────────────────────────────────────────────────────────────
out = r"C:\Users\ocohen\kdesk\provisioning_sf_priority_slide.pptx"
prs.save(out)
print(f"Saved: {out}")
